#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI-OA 业务蓝图 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    sub_frame = sub_box.text_frame
    sub_frame.paragraphs[0].text = subtitle
    sub_frame.paragraphs[0].font.size = Pt(24)
    sub_frame.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
    
    # Table
    cols = len(headers)
    table = slide.shapes.add_table(len(rows)+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4)).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
    
    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 
                   "AI-OA 智能化OA系统",
                   "业务蓝图 | 2026-04-05")
    
    # Slide 2: 目录
    add_content_slide(prs, "目录", [
        "项目概述",
        "核心功能模块",
        "技术架构",
        "AI能力亮点",
        "部署方案",
        "竞争优势分析",
        "实施计划",
        "总结与展望"
    ])
    
    # Slide 3: 项目概述
    add_content_slide(prs, "项目概述", [
        "基于RuoYi开源框架的智能化OA系统",
        "整合OCR识别、n8n流程引擎、AI大模型能力",
        "面向不动产科技(PropTech)赛道",
        "目标用户：中小企业至大型企业",
        "支持私有化部署，灵活扩展",
        "核心理念：AI赋能，提升效率"
    ])
    
    # Slide 4: 核心功能模块
    add_table_slide(prs, "核心功能模块", 
                   ["模块", "功能", "技术亮点"],
                   [
                       ["F1 基础管理", "用户/部门/权限/系统配置", "Spring Security"],
                       ["F2 财务审批", "OCR发票识别+自动填单", "PaddleOCR + 阿里云OCR"],
                       ["F3 智能报表", "周报/月报/年刊自动生成", "GPT-4o + DALL-E 3"],
                       ["F4 AI助手", "RAG知识库+自主学习", "多模型按功能分配"],
                       ["F5 流程中心", "n8n可视化工作流", "Webhook + REST API"],
                       ["F7 企业聊天", "即时消息+群聊+文件", "WebSocket + Kafka"],
                       ["移动端", "鸿蒙+iOS+Android", "原生开发+跨平台适配"]
                   ])
    
    # Slide 5: 技术架构
    add_content_slide(prs, "技术架构", [
        "后端：Java 17+ / Spring Boot 3.x / Spring Cloud",
        "前端：Vue 3 + Element UI",
        "数据库：MySQL 8.0 主备库 + 读写分离",
        "缓存：Redis Cluster (L1+L2多级)",
        "消息队列：Kafka + RabbitMQ 双队列",
        "文件存储：MinIO (S3兼容)",
        "工作流：n8n 可视化流程引擎"
    ])
    
    # Slide 6: AI能力亮点
    add_content_slide(prs, "AI能力亮点", [
        "🤖 多模型按功能分配：GPT-4o/Claude/Kimi 各司其职",
        "📄 OCR 90%+ 准确率：发票/行程单自动识别",
        "🔍 RAG 知识库问答：基于向量检索的智能回答",
        "📊 AI 生成报表：自动生成周报/月报/年刊",
        "🎨 AI 生图/视频：为报表自动配图",
        "📧 智能邮件：OCR疑问单据自动通知审批人",
        "🔗 跳转链接：AI回复包含文档/审批跳转入口"
    ])
    
    # Slide 7: 部署方案
    add_table_slide(prs, "部署方案", 
                   ["方案", "适用场景", "并发用户", "服务器数量"],
                   [
                       ["单体部署", "中小企业", "<100", "1-2台"],
                       ["微服务部署", "中大型企业", "100-500", "8-12台"],
                       ["Docker部署", "开发测试", "<200", "2-4台"],
                       ["Kubernetes", "大型企业", "1000+", "15-20+台"]
                   ])
    
    # Slide 8: 竞争优势
    add_content_slide(prs, "竞争优势分析", [
        "✅ 私有化部署能力强：数据完全自主可控",
        "✅ AI+RAG深度整合：知识库+智能问答差异化",
        "✅ OCR自动化审批：提升财务处理效率80%",
        "✅ 多模型按需配置：成本与效果平衡",
        "✅ 开源为主：降低License成本",
        "✅ 灵活扩展：从单体到K8s随业务成长"
    ])
    
    # Slide 9: 实施计划
    add_table_slide(prs, "实施计划 (Sprint划分)", 
                   ["Sprint", "周期", "目标模块", "交付物"],
                   [
                       ["Sprint 1", "Week 5-7", "F1基础+F4 AI基础", "用户管理+AI模型配置"],
                       ["Sprint 2", "Week 7-9", "F2财务审批", "OCR+审批流程"],
                       ["Sprint 3", "Week 10-12", "F7聊天+F4知识库", "即时消息+RAG问答"],
                       ["Sprint 4", "Week 13-14", "F3/F5/F6增强", "报表+流程+增强功能"]
                   ])
    
    # Slide 10: 总结
    add_content_slide(prs, "总结与展望", [
        "📋 定位清晰：面向PropTech赛道的智能化OA",
        "🏗️ 架构完整：从需求到部署方案全覆盖",
        "🤖 AI差异化：OCR+RAG+多模型整合",
        "📦 交付规范：文档+UI原型+部署脚本",
        "🚀 前景广阔：私有化部署+灵活扩展",
        "💡 持续进化：基于用户反馈不断优化"
    ])
    
    # Slide 11: Thank you
    add_title_slide(prs, 
                   "谢谢！",
                   "AI-OA 智能化OA系统 | 让办公更智能")
    
    # Save
    output_path = "/root/workspace/AI-OA/docs/AI-OA_业务蓝图.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")

if __name__ == "__main__":
    main()
